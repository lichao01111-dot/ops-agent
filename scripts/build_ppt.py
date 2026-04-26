"""Generate JARVIS architecture DETAILED briefing PPT.

Run:
    ./.venv/bin/python scripts/build_ppt.py

Output: docs/ops_agent_architecture_briefing.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------- 常量 ----------

CHINESE_FONT = "PingFang SC"
ENGLISH_FONT = "Helvetica Neue"
MONO_FONT = "Menlo"

# 色板
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
DEEP_BLUE = RGBColor(0x2E, 0x5B, 0xBA)
LIGHT_BLUE = RGBColor(0xE8, 0xF1, 0xFA)
ACCENT = RGBColor(0xF6, 0x92, 0x2A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xA0, 0x6B)
RED = RGBColor(0xD0, 0x45, 0x45)
PURPLE = RGBColor(0x7A, 0x4F, 0xB5)
CODE_BG = RGBColor(0x2B, 0x2B, 0x2B)
CODE_TEXT = RGBColor(0xE0, 0xE0, 0xE0)
CODE_KW = RGBColor(0x7B, 0xB8, 0xF5)
CODE_STR = RGBColor(0xAB, 0xD5, 0x8E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TOTAL_SLIDES = 32


# ---------- 基础工具 ----------


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_font(run, *, size=18, bold=False, color=DARK_GRAY, font=CHINESE_FONT):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide, left, top, width, height, text,
    *, size=18, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
    font=CHINESE_FONT, anchor=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color, font=font)
    return box


def add_multiline(slide, left, top, width, height, lines, *, size=13,
                  color=DARK_GRAY, align=PP_ALIGN.LEFT, bold=False, font=CHINESE_FONT,
                  line_spacing=1.15):
    """Each element of `lines` can be str or (str, kwargs-dict) for per-line override."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, overrides = item
        else:
            text, overrides = item, {}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = overrides.get("align", align)
        p.space_after = Pt(int(size * (line_spacing - 1.0) * 10))
        r = p.add_run()
        r.text = text
        set_font(
            r,
            size=overrides.get("size", size),
            bold=overrides.get("bold", bold),
            color=overrides.get("color", color),
            font=overrides.get("font", font),
        )
    return box


def add_rect(slide, left, top, width, height, *, fill=LIGHT_BLUE, line=None, line_w=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_rect_text(slide, left, top, width, height, text, *,
                  fill=LIGHT_BLUE, text_color=NAVY, size=14, bold=True,
                  line=None, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  font=CHINESE_FONT):
    rect = add_rect(slide, left, top, width, height, fill=fill, line=line)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=text_color, font=font)
    return rect


def add_header(slide, title, subtitle=None, page_num=None, chapter=None):
    # 顶部蓝条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DEEP_BLUE
    bar.line.fill.background()

    # 章节标签
    if chapter:
        add_textbox(
            slide, Inches(0.5), Inches(0.22), Inches(8), Inches(0.3),
            chapter, size=11, bold=True, color=ACCENT,
            font=ENGLISH_FONT,
        )

    # 主标题
    add_textbox(
        slide, Inches(0.5), Inches(0.5), Inches(11), Inches(0.6),
        title, size=26, bold=True, color=NAVY,
    )
    if subtitle:
        add_textbox(
            slide, Inches(0.5), Inches(1.05), Inches(11), Inches(0.4),
            subtitle, size=13, color=MEDIUM_GRAY,
        )

    if page_num is not None:
        add_textbox(
            slide, Inches(12.1), Inches(0.22), Inches(1.0), Inches(0.3),
            f"{page_num:02d} / {TOTAL_SLIDES:02d}",
            size=10, color=MEDIUM_GRAY, align=PP_ALIGN.RIGHT, font=ENGLISH_FONT,
        )


def add_footer(slide):
    add_textbox(
        slide, Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.25),
        "JARVIS · Agent Kernel + Vertical Agent · 详细架构讲解",
        size=9, color=MEDIUM_GRAY, align=PP_ALIGN.RIGHT, font=ENGLISH_FONT,
    )


def add_code_block(slide, left, top, width, height, code_lines, *, size=11):
    """Dark-theme code/text block."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CODE_BG
    rect.line.fill.background()
    tf = rect.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        set_font(r, size=size, color=CODE_TEXT, font=MONO_FONT)
    return rect


def add_arrow(slide, x1, y1, x2, y2, *, color=MEDIUM_GRAY, width=1.5):
    """Line with arrow from (x1,y1) to (x2,y2) in Inches."""
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)  # 1 = STRAIGHT
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    # add arrow end
    from pptx.oxml.ns import qn
    from lxml import etree
    line = connector.line._get_or_add_ln()
    tailEnd = etree.SubElement(line, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("len", "med")
    return connector


# ---------- 页面实现 ----------


def slide_01_cover(prs):
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()

    accent = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), Inches(0.18), Inches(1.8)
    )
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    add_textbox(s, Inches(0.8), Inches(2.5), Inches(12), Inches(1.3),
                "JARVIS 架构详解", size=56, bold=True, color=WHITE)
    add_textbox(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.8),
                "Agent Kernel + Vertical Agent",
                size=30, color=RGBColor(0xCC, 0xDD, 0xEE))
    add_textbox(s, Inches(0.8), Inches(4.4), Inches(12), Inches(0.6),
                "从设计原则到实现细节 · 全栈技术汇报",
                size=18, color=RGBColor(0x99, 0xAA, 0xCC))

    add_textbox(s, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
                "2026 · 技术评审", size=12,
                color=RGBColor(0x88, 0x99, 0xAA), font=ENGLISH_FONT)


def slide_02_toc(prs):
    s = blank_slide(prs)
    add_header(s, "目录", "Table of Contents", page_num=2)

    chapters = [
        ("Part 1", "背景与问题", "为什么要重构", "p.3–5", DEEP_BLUE),
        ("Part 2", "整体分层", "Kernel / Vertical / Supervisor", "p.6–8", GREEN),
        ("Part 3", "Agent Kernel 深度讲解", "每个组件的职责与接口", "p.9–17", PURPLE),
        ("Part 4", "JARVIS 垂直实现", "作为第一个 Vertical 样例", "p.18–22", ACCENT),
        ("Part 5", "关键执行流程", "审批 / 诊断 / 复合请求", "p.23–26", DEEP_BLUE),
        ("Part 6", "插件化与扩展", "10 个插件点 + 新 Vertical 清单", "p.27–28", GREEN),
        ("Part 7", "演进方向与测试", "降级 / 测试 / Supervisor / 路线图", "p.29–32", PURPLE),
    ]
    y = 1.6
    for tag, name, desc, pages, color in chapters:
        add_rect_text(s, Inches(0.6), Inches(y), Inches(1.5), Inches(0.65),
                      tag, fill=color, text_color=WHITE, size=14, bold=True)
        add_textbox(s, Inches(2.3), Inches(y + 0.05), Inches(7), Inches(0.3),
                    name, size=20, bold=True, color=NAVY)
        add_textbox(s, Inches(2.3), Inches(y + 0.4), Inches(7), Inches(0.3),
                    desc, size=13, color=MEDIUM_GRAY)
        add_textbox(s, Inches(11), Inches(y + 0.2), Inches(2), Inches(0.3),
                    pages, size=12, color=MEDIUM_GRAY, align=PP_ALIGN.RIGHT,
                    font=ENGLISH_FONT)
        y += 0.75

    add_footer(s)


# ====================== Part 1 ======================

def slide_03_problem(prs):
    s = blank_slide(prs)
    add_header(s, "旧架构的根本问题",
               "所有逻辑塞进一个 Agent，最终会崩塌在四个方向",
               page_num=3, chapter="PART 1  背景与问题")

    pains = [
        ("🔀", "路由退化", RED,
         "多领域关键词互相冲突。",
         "\"重启\" 既是运维，\"重启会话\" 又是客服 —\n"
         "单 Agent 关键词路由的准确率断崖式下跌。"),
        ("🧨", "安全边界失守", ACCENT,
         "不同业务审批流差异巨大。",
         "\"重启 Pod\"\"发起转账\"\"修改薪资\" 审批人、\n"
         "触发阈值、风险等级完全不同，不能共用一套判断。"),
        ("🧠", "工具选错", PURPLE,
         "LLM 面对 > 20 个工具就开始犯糊涂。",
         "单 Agent 工具数一多，LLM 的工具检索\n"
         "准确率明显退化 — 这是模型能力的天花板。"),
        ("🧴", "记忆串味", DEEP_BLUE,
         "FACTS / HYPOTHESES 只适合诊断-决策-执行。",
         "到了客服、数据场景这几层根本不贴切，\n"
         "强行复用会污染语义。"),
    ]
    col_w, col_h = 6.15, 2.5
    positions = [(0.5, 1.7), (6.75, 1.7), (0.5, 4.3), (6.75, 4.3)]
    for (emoji, title, color, sub, body), (px, py) in zip(pains, positions):
        add_rect(s, Inches(px), Inches(py), Inches(col_w), Inches(col_h),
                 fill=LIGHT_GRAY)
        # 左侧色条
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(px), Inches(py),
                                 Inches(0.15), Inches(col_h))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(s, Inches(px + 0.3), Inches(py + 0.15),
                    Inches(0.7), Inches(0.6), emoji, size=28)
        add_textbox(s, Inches(px + 1.1), Inches(py + 0.15),
                    Inches(4.5), Inches(0.4), title,
                    size=19, bold=True, color=color)
        add_textbox(s, Inches(px + 1.1), Inches(py + 0.7),
                    Inches(4.8), Inches(0.4), sub,
                    size=12, color=MEDIUM_GRAY, bold=True)
        add_textbox(s, Inches(px + 0.3), Inches(py + 1.2),
                    Inches(5.7), Inches(1.2), body,
                    size=12, color=DARK_GRAY)

    add_rect(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4), fill=NAVY)
    add_textbox(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3),
                "结论：窄而深的垂直 Agent 是护城河；一个 Agent 吃天下是陷阱",
                size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_04_split_insight(prs):
    s = blank_slide(prs)
    add_header(s, "核心判断：70% 骨架无关领域，30% 才是业务",
               "这是整个架构重构的出发点",
               page_num=4, chapter="PART 1  背景与问题")

    # 左栏：70%
    add_rect(s, Inches(0.5), Inches(1.6), Inches(6.15), Inches(5.0),
             fill=LIGHT_BLUE, line=DEEP_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.75), Inches(5.8), Inches(0.5),
                "70%  可复用骨架（领域无关）",
                size=20, bold=True, color=DEEP_BLUE)
    left_items = [
        "Planner / PlanStep / advance",
        "LangGraph StateGraph 编排",
        "ToolRegistry + MCP Gateway",
        "DiagnosisExecutor 多假设模式",
        "MemorySchema + RBAC 分层",
        "ApprovalPolicy + AuditLogger",
        "双入口 chat / chat_stream",
        "三级降级（L1/L2/L3）",
    ]
    y = 2.4
    for item in left_items:
        add_textbox(s, Inches(0.9), Inches(y), Inches(5.5), Inches(0.4),
                    f"✓  {item}", size=14, color=NAVY)
        y += 0.48

    # 右栏：30%
    add_rect(s, Inches(6.85), Inches(1.6), Inches(6.0), Inches(5.0),
             fill=RGBColor(0xFF, 0xF2, 0xE6), line=ACCENT)
    add_textbox(s, Inches(7.05), Inches(1.75), Inches(5.7), Inches(0.5),
                "30%  业务逻辑（Ops 特有）",
                size=20, bold=True, color=ACCENT)
    right_items = [
        "tools/*  K8s / Jenkins / Logs 12 个",
        "config/topology.yaml  服务拓扑",
        "router.py  运维关键词映射",
        "extract_namespace / pod / service",
        "_plan_read_only_tool  只读工具编排",
        "_build_pipeline_plan  发布流程",
        "_format_single_read_only_result",
        "_update_memory_from_tool_output",
    ]
    y = 2.4
    for item in right_items:
        add_textbox(s, Inches(7.25), Inches(y), Inches(5.5), Inches(0.4),
                    f"●  {item}", size=14, color=DARK_GRAY)
        y += 0.48

    # 底部结论
    add_rect(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.55), fill=NAVY)
    add_textbox(s, Inches(0.6), Inches(6.88), Inches(12.1), Inches(0.4),
                "把 70% 抽成 Agent Kernel  →  新垂直 Agent 只写 30%",
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_05_principles(prs):
    s = blank_slide(prs)
    add_header(s, "六条设计原则",
               "Kernel 和 Vertical 分界线的裁判",
               page_num=5, chapter="PART 1  背景与问题")

    principles = [
        ("1", "Kernel 零领域知识",
         "任何关键词、工具名、业务 Schema、格式化函数都不许进 Kernel。"),
        ("2", "插件点显式声明",
         "Kernel 通过抽象基类 + 依赖注入暴露扩展点，不通过猜测或反射。"),
        ("3", "契约先于实现",
         "Pydantic / TypedDict 定义 Plan / PlanStep / ToolSpec / MemoryItem，跨 Vertical 共享。"),
        ("4", "安全边界可配置但不可绕过",
         "RBAC、Approval、side_effect 由 Kernel 强制执行，Vertical 只能填配置不能削弱。"),
        ("5", "每个 Vertical 独立实例",
         "各自的工具、路由、记忆语义，不共享运行时状态。"),
        ("6", "Supervisor 是 \"Agent 的 Planner\"",
         "上层调度不关心下层怎么干，只发 sub-plan 给具名的子 Agent。"),
    ]
    y = 1.6
    for num, title, body in principles:
        # 圆形数字
        circle = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.6), Inches(y), Inches(0.75), Inches(0.75)
        )
        circle.fill.solid(); circle.fill.fore_color.rgb = DEEP_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        set_font(r, size=22, bold=True, color=WHITE, font=ENGLISH_FONT)

        # 标题 + 描述
        add_textbox(s, Inches(1.6), Inches(y + 0.02), Inches(11), Inches(0.35),
                    title, size=17, bold=True, color=NAVY)
        add_textbox(s, Inches(1.6), Inches(y + 0.42), Inches(11), Inches(0.4),
                    body, size=12, color=DARK_GRAY)
        y += 0.9

    add_footer(s)


# ====================== Part 2 ======================


def slide_06_layers_overview(prs):
    s = blank_slide(prs)
    add_header(s, "三层总图：Supervisor / Vertical / Kernel",
               "从下向上是依赖关系，从上向下是调度关系",
               page_num=6, chapter="PART 2  整体分层")

    # Supervisor (top)
    add_rect(s, Inches(2.5), Inches(1.5), Inches(8.3), Inches(0.9),
             fill=LIGHT_GRAY, line=MEDIUM_GRAY)
    add_textbox(s, Inches(2.7), Inches(1.55), Inches(8), Inches(0.4),
                "Supervisor（演进方向）", size=16, bold=True, color=MEDIUM_GRAY)
    add_textbox(s, Inches(2.7), Inches(1.95), Inches(8), Inches(0.4),
                "MetaPlanner · PlanStep.execution_target=\"agent:xxx\" · 跨 Agent 的联合审批 / 审计",
                size=11, color=MEDIUM_GRAY)

    # Arrows Supervisor -> Verticals
    for x in [3.0, 5.5, 8.0, 10.5]:
        add_arrow(s, 6.65, 2.5, x, 3.0, color=MEDIUM_GRAY, width=1)

    # 4 Verticals
    verticals = [
        ("JARVIS", "运维\n（已落地）", GREEN),
        ("CsmAgent", "客服\n（未来）", MEDIUM_GRAY),
        ("DataAgent", "数据\n（未来）", MEDIUM_GRAY),
        ("DocAgent", "文档\n（未来）", MEDIUM_GRAY),
    ]
    x = 1.0
    for name, sub, color in verticals:
        add_rect_text(s, Inches(x), Inches(3.0), Inches(2.7), Inches(1.3),
                      f"{name}\n{sub}", fill=color, text_color=WHITE,
                      size=14, bold=True)
        # Arrow Vertical -> Kernel
        add_arrow(s, x + 1.35, 4.4, x + 1.35, 5.0, color=MEDIUM_GRAY, width=1.5)
        x += 3.05

    # Kernel (bottom)
    add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.9),
             fill=LIGHT_BLUE, line=DEEP_BLUE)
    add_textbox(s, Inches(0.7), Inches(5.1), Inches(12), Inches(0.4),
                "🧱  Agent Kernel （领域无关骨架）", size=17, bold=True, color=NAVY)
    components = [
        "Planner", "Router", "Executor", "ToolRegistry", "MCPClient",
        "MemorySchema", "MemoryBackend", "Approval", "Audit", "Session",
    ]
    cx = 0.7
    cy = 5.6
    for i, c in enumerate(components):
        col = i % 5
        row = i // 5
        add_rect_text(s, Inches(cx + col * 2.4), Inches(cy + row * 0.55),
                      Inches(2.3), Inches(0.5),
                      c, fill=WHITE, text_color=NAVY, size=12, bold=False)

    add_footer(s)


def slide_07_boundary_rules(prs):
    s = blank_slide(prs)
    add_header(s, "Kernel vs Vertical：什么归谁？",
               "这条边界是架构评审的准绳 —— 违反它就是反模式",
               page_num=7, chapter="PART 2  整体分层")

    # Kernel 列
    add_rect(s, Inches(0.5), Inches(1.5), Inches(6.15), Inches(5.3),
             fill=LIGHT_BLUE, line=DEEP_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(5.8), Inches(0.5),
                "🧱  Agent Kernel 必须是什么", size=18, bold=True, color=DEEP_BLUE)
    kernel = [
        ("✅", "编排 / 执行 / 审计 / 脱敏 / 降级", "通用能力"),
        ("✅", "抽象基类 ExecutorBase / RouterBase", "插件契约"),
        ("✅", "Pydantic Schema：Plan / PlanStep / ToolSpec", "数据契约"),
        ("✅", "RouteKey / MemoryLayerKey 等可注册字符串", "开放契约"),
        ("❌", "不许出现 \"pod\" \"jenkins\" \"重启\" 等业务词", ""),
        ("❌", "不许硬编码路由、工具名、关键词、拓扑", ""),
        ("❌", "不许导入 agent_ops 任何模块", ""),
    ]
    y = 2.25
    for sign, text, sub in kernel:
        color = GREEN if sign == "✅" else RED
        add_textbox(s, Inches(0.75), Inches(y), Inches(0.4), Inches(0.4),
                    sign, size=16, bold=True, color=color)
        add_textbox(s, Inches(1.2), Inches(y - 0.02), Inches(5.2), Inches(0.4),
                    text, size=13, color=DARK_GRAY, bold=True)
        if sub:
            add_textbox(s, Inches(1.2), Inches(y + 0.35), Inches(5), Inches(0.3),
                        sub, size=10, color=MEDIUM_GRAY)
            y += 0.7
        else:
            y += 0.45

    # Vertical 列
    add_rect(s, Inches(6.85), Inches(1.5), Inches(6.0), Inches(5.3),
             fill=RGBColor(0xFF, 0xF2, 0xE6), line=ACCENT)
    add_textbox(s, Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.5),
                "🎯  Vertical Agent 才能是什么", size=18, bold=True, color=ACCENT)
    vertical = [
        ("✅", "业务工具（K8s、Jenkins、CRM、SQL）", "调用链 endpoints"),
        ("✅", "业务关键词路由（\"重启\"→ mutation）", "Router 子类"),
        ("✅", "风险审批规则（生产 + mutation = 必审批）", "ApprovalPolicy 子类"),
        ("✅", "记忆层定义（observations / hypotheses ...）", "MemorySchema"),
        ("✅", "业务 Prompt 和诊断打分启发式", "_heuristic_score"),
        ("✅", "业务格式化输出（formatters.py）", ""),
        ("✅", "Planner 子类覆写复合拆分规则", "OpsPlanner"),
    ]
    y = 2.25
    for sign, text, sub in vertical:
        add_textbox(s, Inches(7.1), Inches(y), Inches(0.4), Inches(0.4),
                    sign, size=16, bold=True, color=GREEN)
        add_textbox(s, Inches(7.55), Inches(y - 0.02), Inches(5.1), Inches(0.4),
                    text, size=13, color=DARK_GRAY, bold=True)
        if sub:
            add_textbox(s, Inches(7.55), Inches(y + 0.35), Inches(5), Inches(0.3),
                        sub, size=10, color=MEDIUM_GRAY)
            y += 0.7
        else:
            y += 0.45

    add_footer(s)


def slide_08_directory(prs):
    s = blank_slide(prs)
    add_header(s, "代码组织：agent_kernel/ 和 agent_ops/ 两棵树",
               "每个目录都对应架构图上的一个盒子",
               page_num=8, chapter="PART 2  整体分层")

    # Kernel tree
    add_rect(s, Inches(0.5), Inches(1.6), Inches(6.15), Inches(5.3),
             fill=LIGHT_BLUE, line=DEEP_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(0.4),
                "📁  agent_kernel/", size=17, bold=True, color=DEEP_BLUE,
                font=MONO_FONT)
    kernel_tree = [
        "base_agent.py        # BaseAgent + LangGraph 图",
        "planner.py           # Planner + advance / replan",
        "router.py            # RouterBase",
        "executor.py          # ExecutorBase / FunctionExecutor",
        "schemas.py           # Plan / PlanStep / ToolSpec ...",
        "approval.py          # ApprovalPolicy + ApprovalDecision",
        "audit.py             # AuditLogger + sanitizer / sink",
        "session.py           # SessionStore ABC + InMemory",
        "memory/",
        "  schema.py          # MemorySchema (RBAC)",
        "  backend.py         # MemoryBackend ABC",
        "tools/",
        "  registry.py        # ToolRegistry + ToolSpec 索引",
        "  mcp_gateway.py     # MCP Client + register_server",
        "patterns/            # 可选基类库",
        "  multi_hypothesis.py",
        "  approval_gate.py",
    ]
    add_multiline(s, Inches(0.8), Inches(2.1), Inches(5.8), Inches(4.7),
                  kernel_tree, size=10, font=MONO_FONT, color=NAVY,
                  line_spacing=1.25)

    # Ops tree
    add_rect(s, Inches(6.85), Inches(1.6), Inches(6.0), Inches(5.3),
             fill=RGBColor(0xFF, 0xF2, 0xE6), line=ACCENT)
    add_textbox(s, Inches(7.05), Inches(1.7), Inches(5.7), Inches(0.4),
                "📁  agent_ops/", size=17, bold=True, color=ACCENT,
                font=MONO_FONT)
    ops_tree = [
        "agent.py             # JARVIS 装配（代码类名仍为 OpsAgent）",
        "router.py            # IntentRouter (关键词)",
        "planner.py           # OpsPlanner (中文拆分)",
        "risk_policy.py       # OpsApprovalPolicy",
        "memory_schema.py     # OPS_MEMORY_SCHEMA 6 层",
        "extractors.py        # namespace / pod / service 抽取",
        "formatters.py        # 结果格式化",
        "memory_hooks.py      # tool output → memory 规则",
        "topology.py          # ServiceTopology",
        "tool_setup.py        # 12 个 Ops 工具注册入口",
        "schemas.py           # IntentType / AgentRoute / ...",
        "executors/",
        "  knowledge.py       # KnowledgeExecutor",
        "  read_only.py       # ReadOnlyOpsExecutor",
        "  diagnosis.py       # DiagnosisExecutor (多假设)",
        "  mutation.py        # MutationExecutor",
    ]
    add_multiline(s, Inches(7.2), Inches(2.1), Inches(5.7), Inches(4.7),
                  ops_tree, size=10, font=MONO_FONT, color=DARK_GRAY,
                  line_spacing=1.25)

    add_footer(s)


# ====================== Part 3  Kernel ======================


def slide_09_kernel_components(prs):
    s = blank_slide(prs)
    add_header(s, "Kernel 的 11 个核心组件",
               "每个组件都有明确职责边界和可测试的接口",
               page_num=9, chapter="PART 3  Agent Kernel 深度讲解")

    headers_ = ["组件", "职责", "关键接口 / 文件"]
    widths = [2.6, 5.5, 4.2]
    x = 0.5
    for i, h in enumerate(headers_):
        add_rect_text(s, Inches(x), Inches(1.55), Inches(widths[i]), Inches(0.45),
                      h, fill=NAVY, text_color=WHITE, size=13, bold=True)
        x += widths[i]

    rows = [
        ("BaseAgent", "装配 Planner + Executors + Session，暴露 chat / chat_stream",
         "base_agent.py"),
        ("Planner", "生成 Plan、advance / replan、max_iterations 预算",
         "planner.py · initial_plan / advance"),
        ("RouterBase", "把自然语言请求 → RouteDecision",
         "router.py · async route()"),
        ("ExecutorBase", "单个路由的执行器抽象基类",
         "executor.py · async execute(state)"),
        ("ToolRegistry", "本地 + MCP 工具的统一 ToolSpec + retrieve",
         "tools/registry.py"),
        ("MCPClient", "MCP 服务器网关，零代码接入远程工具",
         "tools/mcp_gateway.py"),
        ("MemorySchema", "记忆层定义 + RBAC 白名单",
         "memory/schema.py · assert_can_write"),
        ("MemoryBackend", "共享记忆存储接口（InMemory / Redis / DB）",
         "memory/backend.py"),
        ("ApprovalPolicy", "审批策略接口 + ApprovalReceipt 校验",
         "approval.py · evaluate / resolve_receipt"),
        ("AuditLogger", "审计日志 + sanitizer / sink 扩展钩子",
         "audit.py · log / add_sanitizer / add_sink"),
        ("SessionStore", "会话 + 消息历史 + 记忆条目存储",
         "session.py (ABC + InMemory)"),
    ]
    y = 2.0
    for idx, row in enumerate(rows):
        fill = WHITE if idx % 2 == 0 else LIGHT_GRAY
        x = 0.5
        for i, cell in enumerate(row):
            add_rect_text(s, Inches(x), Inches(y), Inches(widths[i]), Inches(0.42),
                          cell, fill=fill,
                          text_color=NAVY if i == 0 else DARK_GRAY,
                          size=11, bold=(i == 0),
                          align=PP_ALIGN.LEFT if i > 0 else PP_ALIGN.CENTER,
                          font=MONO_FONT if i == 2 else CHINESE_FONT)
            x += widths[i]
        y += 0.42

    add_footer(s)


def slide_10_base_agent_graph(prs):
    s = blank_slide(prs)
    add_header(s, "BaseAgent 的核心图（LangGraph StateGraph）",
               "一图读懂：planner 节点是中心，executors 是动态扇出",
               page_num=10, chapter="PART 3  Agent Kernel 深度讲解")

    # entry
    add_rect_text(s, Inches(5.5), Inches(1.6), Inches(2.3), Inches(0.5),
                  "▶ entry_point", fill=MEDIUM_GRAY, text_color=WHITE, size=12)
    add_arrow(s, 6.65, 2.1, 6.65, 2.45, color=MEDIUM_GRAY, width=2)

    # planner node
    add_rect_text(s, Inches(5.0), Inches(2.45), Inches(3.3), Inches(0.75),
                  "planner  ─  _planner_node",
                  fill=NAVY, text_color=WHITE, size=14, bold=True,
                  font=MONO_FONT)

    # conditional edges down to 4 executor nodes
    executors = [
        ("knowledge", GREEN, 0.7),
        ("read_only_ops", DEEP_BLUE, 3.95),
        ("diagnosis", PURPLE, 7.2),
        ("mutation", ACCENT, 10.45),
    ]
    for name, color, x in executors:
        # arrow
        add_arrow(s, 6.65, 3.25, x + 1.1, 3.9, color=MEDIUM_GRAY, width=1.2)
        # box
        add_rect_text(s, Inches(x), Inches(3.9), Inches(2.2), Inches(0.7),
                      name, fill=color, text_color=WHITE,
                      size=13, bold=True, font=MONO_FONT)
        # arrow back up
        add_arrow(s, x + 1.1, 4.65, 6.65, 5.3, color=RGBColor(0xBB, 0xBB, 0xBB),
                  width=1)

    # finish
    add_rect_text(s, Inches(5.3), Inches(5.3), Inches(2.7), Inches(0.55),
                  "finish  ─  END", fill=RED, text_color=WHITE,
                  size=13, bold=True, font=MONO_FONT)

    # annotations
    add_textbox(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.4),
                "🔑 关键点", size=14, bold=True, color=NAVY)
    anno = [
        "▸ 执行器节点不是写死的 —— 由 Vertical 传入的 executors 列表决定，BaseAgent 动态 add_node",
        "▸ 每个 executor 执行完必须回到 planner 节点，由 advance() 决定 CONTINUE / REPLAN / FINISH",
        "▸ execution_target 字段（§7.2）优先于 route，为未来 Supervisor 派发跨 Agent 预留语义",
    ]
    y = 6.5
    for line in anno:
        add_textbox(s, Inches(0.6), Inches(y), Inches(12.2), Inches(0.3),
                    line, size=11, color=DARK_GRAY)
        y += 0.3


def slide_11_planner(prs):
    s = blank_slide(prs)
    add_header(s, "Planner —— 计划生成 + 前进决策",
               "两大职责：initial_plan() 和 advance()",
               page_num=11, chapter="PART 3  Agent Kernel 深度讲解")

    # initial_plan
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6.15), Inches(5.4),
             fill=LIGHT_BLUE, line=DEEP_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(5.8), Inches(0.4),
                "initial_plan(request) → Plan",
                size=16, bold=True, color=DEEP_BLUE, font=MONO_FONT)
    add_multiline(s, Inches(0.7), Inches(2.1), Inches(5.8), Inches(4.7), [
        "1. self._split_compound(message)",
        "   → Kernel 默认返回 [message]",
        "   → OpsPlanner 覆写用中文关键词拆分",
        "",
        "2. for each segment:",
        "     decision = await router.route(sub_request)",
        "     step = PlanStep(route, intent, goal, ...)",
        "",
        "3. 串成 Plan 带 depends_on 依赖关系",
        "",
        "扩展点：",
        "▸ _split_compound(message) -> list[str]",
        "▸ _dedupe_segments(segments, limit=3)",
    ], size=12, color=DARK_GRAY, font=MONO_FONT, line_spacing=1.2)

    # advance
    add_rect(s, Inches(6.85), Inches(1.55), Inches(6), Inches(5.4),
             fill=RGBColor(0xE5, 0xF4, 0xE9), line=GREEN)
    add_textbox(s, Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.4),
                "advance(plan, last_step) → Decision",
                size=16, bold=True, color=GREEN, font=MONO_FONT)
    add_multiline(s, Inches(7.05), Inches(2.1), Inches(5.7), Inches(4.7), [
        "决策逻辑（有优先级）：",
        "",
        "1. iterations ≥ max_iterations → FINISH",
        "   （硬预算，防 AI 死循环，默认 6）",
        "",
        "2. last_step.status == FAILED → FINISH",
        "   （fail-fast，除非 Vertical 覆写）",
        "",
        "3. 还有 PENDING 步骤 → CONTINUE",
        "   （推进 cursor 到下一个 pending）",
        "",
        "4. _maybe_replan(plan, last) → REPLAN",
        "   （Vertical 覆写追加步骤）",
        "",
        "5. 否则 → FINISH",
    ], size=12, color=DARK_GRAY, font=MONO_FONT, line_spacing=1.2)

    add_footer(s)


def slide_12_router(prs):
    s = blank_slide(prs)
    add_header(s, "Router —— 意图识别与路由决策",
               "从自然语言到 RouteDecision",
               page_num=12, chapter="PART 3  Agent Kernel 深度讲解")

    # Left: Contract
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6.15), Inches(5.4),
             fill=LIGHT_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(5.8), Inches(0.4),
                "契约：RouterBase (ABC)", size=16, bold=True, color=DEEP_BLUE,
                font=MONO_FONT)
    add_code_block(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(2.2), [
        "class RouterBase(ABC):",
        "    @abstractmethod",
        "    async def route(",
        "        self, request: ChatRequest",
        "    ) -> RouteDecision:",
        "        ...",
    ], size=11)
    add_textbox(s, Inches(0.7), Inches(4.55), Inches(5.8), Inches(0.3),
                "RouteDecision 字段：", size=13, bold=True, color=NAVY)
    fields = [
        "• intent: IntentTypeKey (\"query_mysql\" ...)",
        "• route:  RouteKey (\"knowledge\" / \"mutation\" ...)",
        "• risk_level: LOW / MEDIUM / HIGH / CRITICAL",
        "• requires_approval: bool",
        "• rationale: str",
    ]
    y = 4.85
    for f in fields:
        add_textbox(s, Inches(0.9), Inches(y), Inches(5.5), Inches(0.3),
                    f, size=11, color=DARK_GRAY, font=MONO_FONT)
        y += 0.3

    # Right: Ops implementation
    add_rect(s, Inches(6.85), Inches(1.55), Inches(6), Inches(5.4),
             fill=RGBColor(0xFF, 0xF2, 0xE6))
    add_textbox(s, Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.4),
                "Ops 实现：IntentRouter（关键词）",
                size=16, bold=True, color=ACCENT, font=MONO_FONT)
    add_multiline(s, Inches(7.05), Inches(2.2), Inches(5.7), Inches(4.5), [
        "常见映射规则：",
        "",
        "   \"MySQL / 地址 / 密码 / 联系\"",
        "       → knowledge / LOW / no-approval",
        "",
        "   \"pod / 日志 / 状态 / 查\"",
        "       → read_only_ops / LOW",
        "",
        "   \"为什么 / 根因 / 故障 / 异常\"",
        "       → diagnosis / MEDIUM",
        "",
        "   \"重启 / 部署 / 回滚 / 删除\"",
        "       → mutation / HIGH / 必审批",
        "",
        "进阶实现可换成 LLM-based Router",
        "（只要继承 RouterBase 即可）",
    ], size=12, color=DARK_GRAY, font=MONO_FONT, line_spacing=1.15)


def slide_13_executor(prs):
    s = blank_slide(prs)
    add_header(s, "ExecutorBase —— 执行器抽象基类",
               "每个 route 对应一个 Executor；Vertical 自由扩展",
               page_num=13, chapter="PART 3  Agent Kernel 深度讲解")

    add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(1.5),
             fill=LIGHT_BLUE)
    add_code_block(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(1.3), [
        "class ExecutorBase(ABC):",
        "    def __init__(self, *, node_name: str, route_name: str):",
        "        self.node_name = node_name    # LangGraph 节点名",
        "        self.route_name = route_name  # 逻辑路由名",
        "",
        "    @abstractmethod",
        "    async def execute(self, state: dict, event_callback=None) -> dict:",
        "        \"\"\"返回 {final_message, tool_calls, sources, ...}\"\"\"",
    ], size=11)

    # 下方：可选的 patterns
    add_textbox(s, Inches(0.5), Inches(3.2), Inches(12), Inches(0.4),
                "Kernel 内置的可选基类（agent_kernel.patterns/）",
                size=16, bold=True, color=NAVY)

    patterns = [
        ("MultiHypothesisExecutor",
         "多假设并行诊断模式 · 5 stage pipeline",
         "✓  _collect_symptoms\n✓  _generate_hypotheses\n✓  _evidence_args_for\n✓  _score_and_summarize\n✓  _persist",
         PURPLE),
        ("ApprovalGateExecutor",
         "审批闸门执行器 · 先验 receipt 再放行",
         "✓  _execute_approved\n   (只在 approval_policy 通过后才会调用)\n\n✓  _denial_message\n   (可覆写拒绝话术)",
         ACCENT),
    ]
    x = 0.5
    for name, desc, body, color in patterns:
        add_rect_text(s, Inches(x), Inches(3.7), Inches(6.15), Inches(0.55),
                      name, fill=color, text_color=WHITE, size=15, bold=True,
                      font=MONO_FONT)
        add_rect(s, Inches(x), Inches(4.25), Inches(6.15), Inches(2.6),
                 fill=LIGHT_GRAY)
        add_textbox(s, Inches(x + 0.2), Inches(4.35), Inches(5.95), Inches(0.4),
                    desc, size=12, bold=True, color=DARK_GRAY)
        add_textbox(s, Inches(x + 0.2), Inches(4.75), Inches(5.95), Inches(2.0),
                    body, size=12, color=DARK_GRAY, font=MONO_FONT)
        x += 6.3


def slide_14_tools(prs):
    s = blank_slide(prs)
    add_header(s, "ToolRegistry + MCP Gateway",
               "本地工具 + 远程 MCP 工具统一 ToolSpec，统一 retrieve",
               page_num=14, chapter="PART 3  Agent Kernel 深度讲解")

    # ToolSpec
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6.15), Inches(5.4),
             fill=LIGHT_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(5.8), Inches(0.4),
                "ToolSpec：工具元数据契约",
                size=16, bold=True, color=DEEP_BLUE, font=MONO_FONT)
    add_code_block(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(3.2), [
        "class ToolSpec(BaseModel):",
        "    name: str",
        "    description: str",
        "    tags: list[str] = []",
        "    route_affinity: list[str] = []",
        "    side_effect: bool = False",
        "    source: ToolSource  # LOCAL / MCP",
        "    parameters_schema: dict = {}",
    ], size=11)
    add_textbox(s, Inches(0.7), Inches(5.55), Inches(5.8), Inches(0.4),
                "关键字段", size=13, bold=True, color=NAVY)
    add_multiline(s, Inches(0.9), Inches(5.95), Inches(5.5), Inches(1.0), [
        "• side_effect=True ⇒ Kernel 自动要求 approval",
        "• route_affinity ⇒ retrieve 时优先匹配该路由的工具",
        "• source=MCP ⇒ 远端工具，通过 MCPClient 转发",
    ], size=11, color=DARK_GRAY)

    # MCP flow
    add_rect(s, Inches(6.85), Inches(1.55), Inches(6), Inches(5.4),
             fill=RGBColor(0xFF, 0xF2, 0xE6))
    add_textbox(s, Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.4),
                "MCP Gateway：远程工具零代码接入",
                size=16, bold=True, color=ACCENT, font=MONO_FONT)
    add_multiline(s, Inches(7.05), Inches(2.2), Inches(5.7), Inches(4.5), [
        "典型调用链：",
        "",
        "  1. MCPClient.register_server(name, url)",
        "  2. await client.load_tools(name)",
        "  3. for tool_spec in discovered:",
        "       registry.register_mcp(spec, handler)",
        "",
        "4. 后续 registry.retrieve(goal, route, top_k)",
        "   把本地 + MCP 工具一起返回",
        "",
        "5. 调用时 handler.ainvoke(args) 统一接口",
        "   → 本地直接执行",
        "   → MCP 转发到远端服务器",
        "",
        "效果：不写一行代码就能接入一个",
        "      新的 MCP-compatible 工具服务",
    ], size=12, color=DARK_GRAY, font=MONO_FONT, line_spacing=1.15)


def slide_15_memory(prs):
    s = blank_slide(prs)
    add_header(s, "Memory —— 分层记忆 + RBAC + 后端抽象",
               "Schema 管谁能写什么层；Backend 管怎么存",
               page_num=15, chapter="PART 3  Agent Kernel 深度讲解")

    # MemorySchema
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6.15), Inches(5.4),
             fill=LIGHT_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(5.8), Inches(0.4),
                "MemorySchema —— RBAC 层定义",
                size=16, bold=True, color=DEEP_BLUE, font=MONO_FONT)
    add_code_block(s, Inches(0.7), Inches(2.2), Inches(5.8), Inches(1.8), [
        "schema = MemorySchema(layers={",
        "    \"facts\":        {\"knowledge\"},",
        "    \"observations\": {\"read_ops\"},",
        "    \"hypotheses\":   {\"diagnosis\"},",
        "    \"plans\":        {\"planner\"},",
        "})",
    ], size=11)
    add_textbox(s, Inches(0.7), Inches(4.15), Inches(5.8), Inches(0.4),
                "关键 API", size=13, bold=True, color=NAVY)
    add_multiline(s, Inches(0.9), Inches(4.5), Inches(5.6), Inches(2.4), [
        "• assert_can_write(layer, writer)",
        "    非白名单身份 → PermissionError",
        "",
        "• session_store.write_memory_item(...)",
        "    每次写入自动调 assert_can_write",
        "",
        "• session_store.resolve_memory_value(...)",
        "    按优先顺序读取多层",
    ], size=11, color=DARK_GRAY, font=MONO_FONT, line_spacing=1.3)

    # MemoryBackend
    add_rect(s, Inches(6.85), Inches(1.55), Inches(6), Inches(5.4),
             fill=RGBColor(0xE5, 0xF4, 0xE9))
    add_textbox(s, Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.4),
                "MemoryBackend —— 存储可替换",
                size=16, bold=True, color=GREEN, font=MONO_FONT)
    add_multiline(s, Inches(7.05), Inches(2.2), Inches(5.7), Inches(4.5), [
        "class MemoryBackend(ABC):",
        "    def get(session, key): ...",
        "    def put(session, key, val): ...",
        "    def list(session, prefix): ...",
        "",
        "已有实现：",
        "  • InMemoryMemoryBackend (默认)",
        "",
        "未来可替换：",
        "  • RedisMemoryBackend",
        "  • PostgresMemoryBackend",
        "  • VectorMemoryBackend",
        "",
        "⚠  每个 Vertical 持有独立实例",
        "   (§5.5 约束：不共享状态)",
    ], size=12, color=DARK_GRAY, font=MONO_FONT, line_spacing=1.2)


def slide_16_approval(prs):
    s = blank_slide(prs)
    add_header(s, "Approval —— 审批凭据绑定步骤",
               "只认 ApprovalReceipt，不认 context.approved=true",
               page_num=16, chapter="PART 3  Agent Kernel 深度讲解")

    # 核心契约
    add_rect(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(2.0),
             fill=LIGHT_GRAY)
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.4),
                "ApprovalReceipt 的 5 个关键字段",
                size=15, bold=True, color=NAVY)
    add_multiline(s, Inches(0.7), Inches(2.1), Inches(12), Inches(1.4), [
        "• receipt_id     — 唯一凭据 ID，可追溯",
        "• step_id        — 绑定到具体 PlanStep（换步骤就失效）",
        "• approved_by    — 谁批的",
        "• scope          — 生效范围（某个 namespace / 某笔金额）",
        "• expires_at     — 过期时间（默认几分钟）",
    ], size=13, color=DARK_GRAY, font=MONO_FONT, line_spacing=1.3)

    # 验证流程
    add_textbox(s, Inches(0.5), Inches(3.75), Inches(12), Inches(0.4),
                "ApprovalPolicy.evaluate() 验证流程",
                size=16, bold=True, color=NAVY)

    gates = [
        ("①", "step 存在且 requires_approval=True", "否则拒绝"),
        ("②", "context 里有合法的 approval_receipt dict", "否则拒绝"),
        ("③", "receipt.step_id == 当前 step.step_id", "防止凭据复用"),
        ("④", "receipt.expires_at > 现在", "过期即失效"),
        ("⑤", "Vertical 可覆写 validate_receipt 做额外校验", "金额 / namespace"),
    ]
    y = 4.25
    for num, gate, note in gates:
        add_rect_text(s, Inches(0.6), Inches(y), Inches(0.5), Inches(0.45),
                      num, fill=ACCENT, text_color=WHITE, size=14, bold=True)
        add_rect(s, Inches(1.2), Inches(y), Inches(11.6), Inches(0.45),
                 fill=WHITE, line=LIGHT_GRAY)
        add_textbox(s, Inches(1.4), Inches(y + 0.02), Inches(7), Inches(0.4),
                    gate, size=13, color=DARK_GRAY, bold=True)
        add_textbox(s, Inches(8.5), Inches(y + 0.02), Inches(4.2), Inches(0.4),
                    f"失败 → {note}" if "否则" in note else note,
                    size=12, color=MEDIUM_GRAY)
        y += 0.55


def slide_17_audit(prs):
    s = blank_slide(prs)
    add_header(s, "Audit —— 全量可审计 + 脱敏 + SIEM 推送",
               "每次工具调用产生一条 AuditEntry；支持 sanitizer / sink 扩展",
               page_num=17, chapter="PART 3  Agent Kernel 深度讲解")

    # AuditEntry fields
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6.15), Inches(5.4),
             fill=LIGHT_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(5.8), Inches(0.4),
                "AuditEntry 字段", size=16, bold=True, color=DEEP_BLUE,
                font=MONO_FONT)
    fields = [
        "timestamp      datetime",
        "user_id        str",
        "session_id     str",
        "intent         Optional[IntentTypeKey]",
        "route          Optional[RouteKey]",
        "risk_level     Optional[RiskLevel]",
        "needs_approval bool",
        "tool_name      Optional[str]",
        "tool_calls     list[str]",
        "params         dict      # ← sanitize 目标",
        "result_summary str",
        "success        bool",
        "duration_ms    int",
    ]
    y = 2.15
    for f in fields:
        add_textbox(s, Inches(0.85), Inches(y), Inches(5.6), Inches(0.32),
                    f, size=11, color=DARK_GRAY, font=MONO_FONT)
        y += 0.35

    # Sanitizer + Sink
    add_rect(s, Inches(6.85), Inches(1.55), Inches(6), Inches(5.4),
             fill=RGBColor(0xE5, 0xF4, 0xE9))
    add_textbox(s, Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.4),
                "扩展钩子：sanitizer / sink",
                size=16, bold=True, color=GREEN, font=MONO_FONT)
    add_code_block(s, Inches(7.05), Inches(2.2), Inches(5.7), Inches(2.4), [
        "logger.add_sanitizer(lambda params: {",
        "    **params,",
        "    'password': '***',",
        "    'token':    '***',",
        "})",
        "",
        "logger.add_sink(siem_sink)  # 写 SIEM",
        "logger.add_sink(metrics_sink) # 上报指标",
    ], size=11)
    add_multiline(s, Inches(7.05), Inches(4.75), Inches(5.7), Inches(2.2), [
        "▸ Sanitizers 顺序执行，失败不影响下游",
        "▸ Sinks 独立异常隔离（一个挂了不传染）",
        "▸ 默认 sanitizer 已覆盖 password/token/ak/sk",
        "",
        "▸ Kernel 强制：所有 tool call 必过 log()",
        "  →  Ops / 客服 / 数据 Agent 自动同样受审计",
    ], size=12, color=DARK_GRAY, font=CHINESE_FONT, line_spacing=1.2)


def slide_18_invariants(prs):
    s = blank_slide(prs)
    add_header(s, "Kernel 的 5 条不变量（§4.2）",
               "任何 Vertical 都不能违反 —— 业务代码绕不过去",
               page_num=18, chapter="PART 3  Agent Kernel 深度讲解")

    rules = [
        ("#1", "side_effect 工具 + 必须 receipt",
         "side_effect=True 的工具只能被 requires_approval=True 的 PlanStep 调用，\n"
         "且必须携带与该 step 绑定、可校验、未过期的 approval_receipt",
         "E2E 测试 B01-B04 覆盖", RED),
        ("#2", "所有工具调用走 _invoke_tool",
         "必审计、必脱敏；业务代码不能越过 Kernel 直接调 handler",
         "E2E 测试 B08 覆盖", ACCENT),
        ("#3", "记忆写入必须过 RBAC",
         "write_memory_item 先 assert_can_write；非法 writer → PermissionError",
         "E2E 测试 B05 覆盖", PURPLE),
        ("#4", "max_iterations 是硬预算",
         "超过立即 FINISH，防止 AI 死循环烧 Token",
         "E2E 测试 B06 覆盖", DEEP_BLUE),
        ("#5", "FAILED 默认 fail-fast",
         "一步失败就停车；Vertical 可通过 _maybe_replan 覆写，但必须显式",
         "E2E 测试 B07 覆盖", GREEN),
    ]
    y = 1.6
    for tag, title, body, test, color in rules:
        add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.0),
                 fill=LIGHT_GRAY)
        add_rect(s, Inches(0.5), Inches(y), Inches(0.15), Inches(1.0), fill=color)
        add_textbox(s, Inches(0.8), Inches(y + 0.05), Inches(1.2), Inches(0.4),
                    tag, size=18, bold=True, color=color, font=ENGLISH_FONT)
        add_textbox(s, Inches(2.0), Inches(y + 0.05), Inches(9), Inches(0.4),
                    title, size=15, bold=True, color=NAVY)
        add_textbox(s, Inches(2.0), Inches(y + 0.48), Inches(9), Inches(0.5),
                    body, size=11, color=DARK_GRAY)
        add_textbox(s, Inches(11.2), Inches(y + 0.35), Inches(1.7), Inches(0.4),
                    test, size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
        y += 1.08


# ====================== Part 4  JARVIS ======================


def slide_19_ops_composition(prs):
    s = blank_slide(prs)
    add_header(s, "JARVIS 组成：装配就是套壳",
               "create_ops_agent() 把 Kernel 组件拼起来",
               page_num=19, chapter="PART 4  JARVIS 垂直实现")

    add_code_block(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(4.6), [
        "def create_ops_agent() -> OpsAgent:",
        "    registry      = create_tool_registry()",
        "    register_ops_builtins(registry)              # 12 个 Ops 工具注册",
        "",
        "    audit_logger  = create_audit_logger()        # 默认脱敏已内置",
        "    session_store = create_session_store(",
        "        memory_schema=OPS_MEMORY_SCHEMA,         # 6 层 Ops 语义",
        "    )",
        "    mcp_client    = create_mcp_client(registry=registry)",
        "",
        "    return OpsAgent(",
        "        session_store   = session_store,",
        "        tool_registry   = registry,",
        "        audit_logger    = audit_logger,",
        "        mcp_client      = mcp_client,",
        "    )",
        "",
        "# OpsAgent.__init__ 内部：",
        "#   router          = IntentRouter()             # 关键词路由",
        "#   planner         = OpsPlanner(router=router)  # 中文复合拆分",
        "#   approval_policy = OpsApprovalPolicy()        # 生产 namespace 必审批",
        "#   executors       = [Knowledge, ReadOnly, Diagnosis, Mutation]",
    ], size=12)

    add_rect(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.5), fill=NAVY)
    add_textbox(s, Inches(0.6), Inches(6.38), Inches(12.1), Inches(0.35),
                "想做新 Vertical？复制这个文件，替换 4 个 Ops 专有组件即可",
                size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_20_ops_executors(prs):
    s = blank_slide(prs)
    add_header(s, "JARVIS 的 4 个 Executor",
               "每个路由都有明确的职责、工具范围、记忆写入权限",
               page_num=20, chapter="PART 4  JARVIS 垂直实现")

    executors = [
        ("KnowledgeExecutor", GREEN,
         "route = knowledge",
         "回答「MySQL 地址」「联系人」等静态知识",
         "query_knowledge",
         "facts 层"),
        ("ReadOnlyOpsExecutor", DEEP_BLUE,
         "route = read_only_ops",
         "查 pod / 日志 / 部署状态 / Jenkins 构建",
         "get_pod_status, get_pod_logs,\nsearch_logs, query_jenkins_build ...",
         "observations 层"),
        ("DiagnosisExecutor ★", PURPLE,
         "route = diagnosis",
         "多假设并行 + 证据收集 + 打分 + 根因总结",
         "diagnose_pod, get_pod_status,\nsearch_logs, get_error_statistics ...",
         "hypotheses 层"),
        ("MutationExecutor", ACCENT,
         "route = mutation",
         "变更操作：重启 / 发布 / 回滚，必带 receipt",
         "restart_pod, restart_deployment,\ntrigger_jenkins_build ...",
         "plans + execution 层"),
    ]
    col_w = 2.95
    y_top = 1.6
    x = 0.5
    for name, color, route, desc, tools, mem in executors:
        # header
        add_rect_text(s, Inches(x), Inches(y_top), Inches(col_w), Inches(0.55),
                      name, fill=color, text_color=WHITE, size=13, bold=True,
                      font=MONO_FONT)
        # body
        add_rect(s, Inches(x), Inches(y_top + 0.55), Inches(col_w), Inches(5),
                 fill=LIGHT_GRAY)
        add_textbox(s, Inches(x + 0.1), Inches(y_top + 0.7), Inches(col_w - 0.2),
                    Inches(0.35), route, size=11, color=color, bold=True,
                    font=MONO_FONT)
        add_textbox(s, Inches(x + 0.1), Inches(y_top + 1.1), Inches(col_w - 0.2),
                    Inches(0.4), "职责", size=11, bold=True, color=NAVY)
        add_textbox(s, Inches(x + 0.1), Inches(y_top + 1.4), Inches(col_w - 0.2),
                    Inches(1.0), desc, size=11, color=DARK_GRAY)

        add_textbox(s, Inches(x + 0.1), Inches(y_top + 2.6), Inches(col_w - 0.2),
                    Inches(0.4), "典型工具", size=11, bold=True, color=NAVY)
        add_textbox(s, Inches(x + 0.1), Inches(y_top + 2.9), Inches(col_w - 0.2),
                    Inches(1.8), tools, size=10, color=DARK_GRAY, font=MONO_FONT)

        add_textbox(s, Inches(x + 0.1), Inches(y_top + 4.75), Inches(col_w - 0.2),
                    Inches(0.4), "记忆写入", size=11, bold=True, color=NAVY)
        add_textbox(s, Inches(x + 0.1), Inches(y_top + 5.05), Inches(col_w - 0.2),
                    Inches(0.4), mem, size=11, color=color, bold=True)
        x += col_w + 0.1

    add_footer(s)


def slide_21_ops_planner(prs):
    s = blank_slide(prs)
    add_header(s, "OpsPlanner —— 复合请求的中文拆分",
               "_split_compound 覆写示例：把一句话拆成多步",
               page_num=21, chapter="PART 4  JARVIS 垂直实现")

    # Example flow
    add_textbox(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.4),
                "用户原句", size=14, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.5),
             fill=LIGHT_BLUE, line=DEEP_BLUE)
    add_textbox(s, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.4),
                "「先查一下 staging pod 状态，然后帮我重启 order-service」",
                size=16, color=DEEP_BLUE, bold=True, align=PP_ALIGN.CENTER)

    # 箭头
    add_arrow(s, 6.65, 2.55, 6.65, 2.9, color=MEDIUM_GRAY, width=2)

    # 拆分规则
    add_textbox(s, Inches(0.5), Inches(3.0), Inches(12), Inches(0.4),
                "OpsPlanner._split_compound() 的正则规则", size=14, bold=True, color=NAVY)
    add_code_block(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.5), [
        "_OPS_SPLIT_PATTERNS = [",
        "    re.compile(r\"\\s*然后\\s*\"),",
        "    re.compile(r\"\\s*接着\\s*\"),",
        "    re.compile(r\"\\s*再\\s*(?=(?:帮|把|重|触|回|生|执))\"),",
        "    re.compile(r\"\\s*,\\s*然后\\s*\"),  # 中英逗号两种",
        "]",
    ], size=11)

    add_arrow(s, 6.65, 4.95, 6.65, 5.3, color=MEDIUM_GRAY, width=2)

    # 拆分结果
    add_textbox(s, Inches(0.5), Inches(5.4), Inches(12), Inches(0.4),
                "拆出 2 个 PlanStep（带依赖关系）", size=14, bold=True, color=NAVY)
    add_rect_text(s, Inches(0.5), Inches(5.85), Inches(6.1), Inches(1.15),
                  "Step 1\n查 staging pod 状态\nroute = read_only_ops · LOW",
                  fill=DEEP_BLUE, text_color=WHITE, size=13, bold=True)
    # 依赖箭头
    add_arrow(s, 6.6, 6.42, 6.9, 6.42, color=ACCENT, width=2)
    add_rect_text(s, Inches(6.9), Inches(5.85), Inches(5.9), Inches(1.15),
                  "Step 2  (depends_on=[Step1])\n重启 order-service\nroute = mutation · HIGH · 必审批",
                  fill=ACCENT, text_color=WHITE, size=13, bold=True)


def slide_22_ops_approval_memory(prs):
    s = blank_slide(prs)
    add_header(s, "OpsApprovalPolicy + OPS_MEMORY_SCHEMA",
               "Ops 填好的两个安全相关插件槽",
               page_num=22, chapter="PART 4  JARVIS 垂直实现")

    # OpsApprovalPolicy
    add_rect(s, Inches(0.5), Inches(1.55), Inches(6.15), Inches(5.4),
             fill=RGBColor(0xFF, 0xF2, 0xE6))
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(5.8), Inches(0.4),
                "OpsApprovalPolicy（risk_policy.py）",
                size=15, bold=True, color=ACCENT, font=MONO_FONT)
    add_textbox(s, Inches(0.7), Inches(2.1), Inches(5.8), Inches(0.4),
                "风险矩阵：", size=13, bold=True, color=NAVY)
    matrix_header = ["namespace", "route", "决定"]
    widths_m = [1.9, 1.9, 2.0]
    y = 2.5
    x = 0.7
    for i, h in enumerate(matrix_header):
        add_rect_text(s, Inches(x), Inches(y), Inches(widths_m[i]), Inches(0.35),
                      h, fill=NAVY, text_color=WHITE, size=11, bold=True)
        x += widths_m[i]
    rows_m = [
        ("default / staging", "read_only_ops", "✓ 放行", GREEN),
        ("default / staging", "diagnosis", "✓ 放行", GREEN),
        ("default / staging", "mutation", "⚠ 需 receipt", ACCENT),
        ("production", "read_only_ops", "✓ 放行", GREEN),
        ("production", "mutation", "🛑 必须审批", RED),
        ("production", "delete*", "🛑 双人复核", RED),
    ]
    y += 0.38
    for row in rows_m:
        x = 0.7
        for i, cell in enumerate(row[:3]):
            add_rect_text(s, Inches(x), Inches(y), Inches(widths_m[i]), Inches(0.38),
                          cell, fill=WHITE,
                          text_color=row[3] if i == 2 else DARK_GRAY,
                          size=11, bold=(i == 2),
                          line=LIGHT_GRAY)
            x += widths_m[i]
        y += 0.38

    # OPS_MEMORY_SCHEMA
    add_rect(s, Inches(6.85), Inches(1.55), Inches(6), Inches(5.4),
             fill=LIGHT_BLUE)
    add_textbox(s, Inches(7.05), Inches(1.65), Inches(5.7), Inches(0.4),
                "OPS_MEMORY_SCHEMA (memory_schema.py)",
                size=15, bold=True, color=DEEP_BLUE, font=MONO_FONT)
    add_textbox(s, Inches(7.05), Inches(2.1), Inches(5.7), Inches(0.4),
                "6 层记忆 + 对应 writer", size=13, bold=True, color=NAVY)
    layers = [
        ("facts", "事实层：服务地址、联系人", "knowledge"),
        ("observations", "观察层：pod 状态、日志片段", "read_ops"),
        ("hypotheses", "假设层：诊断假设 + 评分", "diagnosis"),
        ("plans", "计划层：变更计划", "change_planner"),
        ("execution", "执行层：已执行动作", "change_executor"),
        ("verification", "验证层：验证结果", "verifier"),
    ]
    y = 2.55
    for layer, desc, writer in layers:
        add_rect_text(s, Inches(7.05), Inches(y), Inches(1.6), Inches(0.5),
                      layer, fill=DEEP_BLUE, text_color=WHITE, size=11, bold=True,
                      font=MONO_FONT)
        add_textbox(s, Inches(8.75), Inches(y + 0.02), Inches(3.1), Inches(0.5),
                    desc, size=10, color=DARK_GRAY)
        add_textbox(s, Inches(11.85), Inches(y + 0.05), Inches(1.0), Inches(0.4),
                    writer, size=10, color=ACCENT, bold=True, font=MONO_FONT)
        y += 0.55


# ====================== Part 5  Flow ======================


def slide_23_chat_flow(prs):
    s = blank_slide(prs)
    add_header(s, "一次 chat() 请求的完整流水线",
               "从 HTTP 入口到 Audit 落盘的 9 个阶段",
               page_num=23, chapter="PART 5  关键执行流程")

    stages = [
        ("①", "API 入口", "ChatRequest(message, user_id, session_id)"),
        ("②", "状态构建", "_build_initial_state 读取最近 6 条消息"),
        ("③", "Planner 节点", "initial_plan → 拆分 → 第一个 PlanStep"),
        ("④", "Dispatcher", "按 execution_target 或 route → 挑选节点"),
        ("⑤", "Executor 执行", "执行器调 _invoke_tool 执行工具"),
        ("⑥", "Approval 闸门", "side_effect 工具必须过 ApprovalPolicy.evaluate"),
        ("⑦", "_invoke_tool", "调用 handler.ainvoke + Audit 落盘"),
        ("⑧", "Memory 写入", "按 Schema RBAC 写入对应层"),
        ("⑨", "回到 Planner", "advance → CONTINUE / REPLAN / FINISH"),
    ]
    # 纵向长表
    y = 1.55
    for num, title, body in stages:
        # 数字圆
        c = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.5), Inches(0.5)
        )
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.line.fill.background()
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        set_font(r, size=14, bold=True, color=WHITE)
        # 内容
        add_rect(s, Inches(1.4), Inches(y), Inches(11.4), Inches(0.5),
                 fill=LIGHT_GRAY)
        add_textbox(s, Inches(1.6), Inches(y + 0.05), Inches(3), Inches(0.4),
                    title, size=13, bold=True, color=NAVY)
        add_textbox(s, Inches(4.7), Inches(y + 0.05), Inches(8.0), Inches(0.4),
                    body, size=11, color=DARK_GRAY, font=MONO_FONT)
        y += 0.58

    # 底部注解
    add_textbox(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.3),
                "🔁 ③–⑨ 循环直到 FINISH 或 max_iterations",
                size=12, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER)


def slide_24_approval_gate(prs):
    s = blank_slide(prs)
    add_header(s, "审批闸门的工作原理",
               "side_effect 工具从 _invoke_tool 到 handler 之间必经这个门",
               page_num=24, chapter="PART 5  关键执行流程")

    # flow
    # Step 1
    add_rect_text(s, Inches(0.5), Inches(1.6), Inches(2.5), Inches(0.7),
                  "Executor.execute()", fill=DEEP_BLUE, text_color=WHITE,
                  size=12, bold=True, font=MONO_FONT)
    add_arrow(s, 3.05, 1.95, 3.45, 1.95)
    add_rect_text(s, Inches(3.5), Inches(1.6), Inches(2.8), Inches(0.7),
                  "await _invoke_tool(...)", fill=NAVY, text_color=WHITE,
                  size=12, bold=True, font=MONO_FONT)
    add_arrow(s, 6.35, 1.95, 6.75, 1.95)
    add_rect_text(s, Inches(6.8), Inches(1.6), Inches(3.2), Inches(0.7),
                  "ToolRegistry.get_spec(tool)",
                  fill=LIGHT_GRAY, text_color=DARK_GRAY, size=11,
                  font=MONO_FONT, bold=False)
    add_arrow(s, 10.05, 1.95, 10.45, 1.95)
    add_rect_text(s, Inches(10.5), Inches(1.6), Inches(2.3), Inches(0.7),
                  "spec.side_effect?",
                  fill=ACCENT, text_color=WHITE, size=12, bold=True)

    # 分叉
    # 向下：False
    add_arrow(s, 11.65, 2.35, 11.65, 3.1, color=GREEN)
    add_textbox(s, Inches(11.75), Inches(2.5), Inches(1.3), Inches(0.3),
                "No", size=12, bold=True, color=GREEN)
    add_rect_text(s, Inches(10.5), Inches(3.1), Inches(2.3), Inches(0.6),
                  "直接执行", fill=GREEN, text_color=WHITE, size=12, bold=True)

    # Yes 向左 → ApprovalPolicy
    add_arrow(s, 10.5, 1.95, 6.8, 3.15, color=RED, width=2)
    add_textbox(s, Inches(9.7), Inches(2.3), Inches(1.2), Inches(0.3),
                "Yes", size=12, bold=True, color=RED)
    add_rect_text(s, Inches(4.0), Inches(3.1), Inches(2.8), Inches(0.7),
                  "approval_policy.evaluate()",
                  fill=RED, text_color=WHITE, size=12, bold=True,
                  font=MONO_FONT)

    # evaluate → decision
    add_arrow(s, 5.4, 3.8, 5.4, 4.3, color=MEDIUM_GRAY, width=2)

    # Not approved → FAIL
    add_rect_text(s, Inches(0.7), Inches(4.3), Inches(3.8), Inches(1.5),
                  "❌ 未批准\n\nToolCallEvent.status = FAILED\n"
                  "返回 {\"error\": \"需要审批...\"}\nAudit 照样记录",
                  fill=RED, text_color=WHITE, size=12, bold=True)

    # Approved → run handler
    add_rect_text(s, Inches(6.0), Inches(4.3), Inches(3.8), Inches(1.5),
                  "✅ 批准\n\nhandler.ainvoke(args)\n"
                  "status = SUCCESS\nMemory 写入 + Audit 落盘",
                  fill=GREEN, text_color=WHITE, size=12, bold=True)

    # legend
    add_rect(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
             fill=LIGHT_BLUE)
    add_textbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.3),
                "💡 关键：这条路径由 Kernel 在 _invoke_tool 里强制执行",
                size=13, bold=True, color=NAVY)
    add_textbox(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
                "Vertical 想绕过？唯一办法是不走 _invoke_tool —— 但那就没有审计了，两难。",
                size=11, color=DARK_GRAY)


def slide_25_diagnosis(prs):
    s = blank_slide(prs)
    add_header(s, "DiagnosisExecutor —— 多假设并行诊断",
               "「多个假设 · 并行取证 · 启发式打分 · 归纳结论」",
               page_num=25, chapter="PART 5  关键执行流程")

    stages = [
        ("1", "症状采集", "_collect_symptoms",
         "diagnose_pod\nget_pod_status\nsearch_logs",
         LIGHT_BLUE),
        ("2", "假设生成", "_generate_hypotheses",
         "LLM + 拓扑 + 候选工具\n→ 至多 4 条互不重复\n  Hypothesis",
         RGBColor(0xE5, 0xF4, 0xE9)),
        ("3", "并行取证", "_collect_evidence_parallel",
         "asyncio.gather\n每假设 ≤ 2 个\n证据工具",
         RGBColor(0xFF, 0xF2, 0xE6)),
        ("4", "打分合成", "_score_and_synthesize",
         "启发式：error/oom/\ncrashloop → +1.8\n疑点对象匹配 → +0.5",
         RGBColor(0xF0, 0xE6, 0xFF)),
        ("5", "写入记忆", "_write_memory",
         "每条 hypothesis 一条\n+ top_hypothesis_id\n+ diagnosis_summary",
         LIGHT_BLUE),
    ]
    x = 0.4
    col_w = 2.55
    y_top = 1.6
    for (num, title, func, body, color), i in zip(stages, range(len(stages))):
        # 序号
        c = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + col_w/2 - 0.25), Inches(y_top),
            Inches(0.5), Inches(0.5)
        )
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
        c.line.fill.background()
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        set_font(r, size=14, bold=True, color=WHITE)

        # 卡片
        add_rect(s, Inches(x), Inches(y_top + 0.6), Inches(col_w), Inches(3.2),
                 fill=color)
        add_textbox(s, Inches(x), Inches(y_top + 0.75), Inches(col_w),
                    Inches(0.4), title, size=15, bold=True, color=NAVY,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(x), Inches(y_top + 1.2), Inches(col_w),
                    Inches(0.35), func, size=11, color=DEEP_BLUE,
                    align=PP_ALIGN.CENTER, font=MONO_FONT)
        add_textbox(s, Inches(x + 0.2), Inches(y_top + 1.7), Inches(col_w - 0.4),
                    Inches(2.0), body, size=11, color=DARK_GRAY,
                    font=MONO_FONT, align=PP_ALIGN.CENTER)

        # arrow to next
        if i < len(stages) - 1:
            add_arrow(s, x + col_w - 0.05, y_top + 2.2, x + col_w + 0.08,
                      y_top + 2.2, color=MEDIUM_GRAY, width=1.5)
        x += col_w + 0.05

    # Fallback path
    add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
             fill=RGBColor(0xFF, 0xE8, 0xD8))
    add_textbox(s, Inches(0.6), Inches(5.08), Inches(12.1), Inches(0.35),
                "🛟 降级路径：假设生成失败 / LLM 不可用 → _fallback_single_chain 仍然返回已收集症状",
                size=12, color=ACCENT, bold=True)

    # Kernel pattern note
    add_rect(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4),
             fill=LIGHT_BLUE)
    add_textbox(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.4),
                "✨ 这个 5-stage pipeline 已被抽成 Kernel 的 MultiHypothesisExecutor 基类",
                size=13, bold=True, color=NAVY)
    add_textbox(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.8),
                "未来客服 Agent 也可以继承：\"退款 / 漏发 / 延迟\" 三条假设并行查 → 打分归纳根因\n"
                "这就是 §5.3 / §6 #10 要求的\"可选基类\"模式库",
                size=11, color=DARK_GRAY)


def slide_26_compound_example(prs):
    s = blank_slide(prs)
    add_header(s, "实战例子：复合请求的全链路时序",
               "「先查 staging pod，然后重启 order-service」",
               page_num=26, chapter="PART 5  关键执行流程")

    add_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.45), fill=NAVY)
    add_textbox(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.35),
                "t=0ms  用户：\"先查一下 staging pod 状态，然后帮我重启 order-service\"",
                size=13, bold=True, color=WHITE)

    events = [
        ("t=5ms", "Planner", "OpsPlanner._split_compound 分出 2 段",
         DEEP_BLUE, "planner"),
        ("t=8ms", "Router × 2", "route(Step1)=read_only_ops / LOW\n"
         "route(Step2)=mutation / HIGH / 必审批",
         GREEN, "router"),
        ("t=30ms", "Step 1 执行", "ReadOnlyOpsExecutor.execute\n"
         "  → _invoke_tool('get_pod_status', namespace=staging)\n"
         "  → Audit 落盘 + Memory 写 observations",
         DEEP_BLUE, "executor"),
        ("t=45ms", "advance()", "Plan cursor → Step 2；返回 CONTINUE",
         MEDIUM_GRAY, "planner"),
        ("t=48ms", "Step 2 进入 Mutation", "Executor 调 _invoke_tool('restart_deployment')\n"
         "  → ToolSpec.side_effect=True\n"
         "  → ApprovalPolicy.evaluate 拒绝 (无 receipt)",
         RED, "approval"),
        ("t=50ms", "返回用户", "\"此操作需要审批，请先批准……\"\n"
         "Audit 照样记录一条 FAILED 条目",
         ACCENT, "finish"),
    ]
    y = 2.1
    for (t, phase, desc, color, who) in events:
        # timestamp
        add_textbox(s, Inches(0.5), Inches(y + 0.05), Inches(1.0), Inches(0.3),
                    t, size=11, color=MEDIUM_GRAY, font=MONO_FONT, bold=True)
        # vertical dot
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.55), Inches(y + 0.12), Inches(0.2), Inches(0.2)
        )
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        # phase
        add_rect_text(s, Inches(1.9), Inches(y), Inches(2.4), Inches(0.5),
                      phase, fill=color, text_color=WHITE, size=12, bold=True)
        # description
        add_rect(s, Inches(4.4), Inches(y), Inches(8.4), Inches(0.8),
                 fill=LIGHT_GRAY)
        add_textbox(s, Inches(4.55), Inches(y + 0.05), Inches(8.2), Inches(0.7),
                    desc, size=10, color=DARK_GRAY, font=MONO_FONT)
        y += 0.82

    add_footer(s)


# ====================== Part 6  Plugins ======================


def slide_27_plugin_points(prs):
    s = blank_slide(prs)
    add_header(s, "10 个插件点：Kernel 对外的所有扩展面",
               "Vertical 就是\"填这些槽位\"",
               page_num=27, chapter="PART 6  插件化与扩展")

    headers = ["#", "插件点", "基类 / 契约", "Ops 填了什么"]
    widths = [0.6, 2.5, 4.5, 5.7]
    x = 0.5
    for i, h in enumerate(headers):
        add_rect_text(s, Inches(x), Inches(1.55), Inches(widths[i]), Inches(0.4),
                      h, fill=NAVY, text_color=WHITE, size=12, bold=True)
        x += widths[i]

    rows = [
        ("1", "路由器", "RouterBase.route() -> RouteDecision",
         "IntentRouter 关键词映射"),
        ("2", "执行器", "ExecutorBase.execute(state) -> dict",
         "Knowledge / ReadOnly / Diagnosis / Mutation"),
        ("3", "工具", "@tool + ToolRegistry.register_local/_mcp",
         "K8s / Jenkins / Logs / Knowledge 共 12 个"),
        ("4", "MCP 服务器", "MCPClient.register_server(name, url)",
         "可接入任意 MCP-compatible 远端工具"),
        ("5", "Planner 定制", "Planner 子类 _split_compound / _maybe_replan",
         "OpsPlanner 中文复合拆分"),
        ("6", "记忆 Schema", "MemorySchema(layers={...})",
         "OPS_MEMORY_SCHEMA 6 层"),
        ("7", "审批策略", "ApprovalPolicy.evaluate(step, context)",
         "OpsApprovalPolicy 风险矩阵"),
        ("8", "审计扩展", "AuditLogger.add_sanitizer / add_sink",
         "Ops 级别脱敏 + SIEM 可扩展"),
        ("9", "RBAC 身份", "AgentIdentityKey 可注册字符串",
         "knowledge / read_ops / diagnosis 等"),
        ("10", "Executor 模式库", "MultiHypothesisExecutor / ApprovalGateExecutor",
         "Ops DiagnosisExecutor 可继承"),
    ]
    y = 1.95
    for idx, row in enumerate(rows):
        fill = WHITE if idx % 2 == 0 else LIGHT_GRAY
        x = 0.5
        for i, cell in enumerate(row):
            add_rect_text(s, Inches(x), Inches(y), Inches(widths[i]),
                          Inches(0.44), cell, fill=fill,
                          text_color=ACCENT if i == 0 else NAVY if i == 1 else DARK_GRAY,
                          size=10 if i > 1 else 11,
                          bold=(i <= 1),
                          align=PP_ALIGN.CENTER if i <= 1 else PP_ALIGN.LEFT,
                          font=MONO_FONT if i == 2 else CHINESE_FONT)
            x += widths[i]
        y += 0.44

    add_footer(s)


def slide_28_new_vertical(prs):
    s = blank_slide(prs)
    add_header(s, "做一个新 Vertical 需要几步？",
               "以假想的「CsmAgent 客服」为例 · ~ 1-2 周工作量",
               page_num=28, chapter="PART 6  插件化与扩展")

    steps = [
        ("1", "定义记忆 Schema",
         "agent_csm/memory_schema.py",
         "CSM_MEMORY_SCHEMA = MemorySchema(layers={\n"
         "    'user_profile':    {'crm_reader'},\n"
         "    'conversation':    {'dialogue'},\n"
         "    'order_context':   {'crm_reader'},\n"
         "    'escalation_plan': {'supervisor'},\n"
         "})"),
        ("2", "定义风险策略",
         "agent_csm/risk_policy.py",
         "class CsmApprovalPolicy(ApprovalPolicy):\n"
         "    def validate_receipt(...):\n"
         "        # 退款 > 1000 元 → 需要主管 receipt"),
        ("3", "定义路由器",
         "agent_csm/router.py",
         "class CsmKeywordRouter(RouterBase):\n"
         "    async def route(request):\n"
         "        if '退款' in msg:  return RouteDecision(..., 'refund', HIGH)\n"
         "        if '物流' in msg:  return RouteDecision(..., 'tracking', LOW)"),
        ("4", "实现执行器 + 注册工具",
         "agent_csm/executors/ 和 tools/",
         "RefundExecutor / TrackingExecutor / EscalationExecutor\n"
         "接入 CRM / 订单系统 / 工单系统"),
        ("5", "装配入口",
         "agent_csm/__init__.py · create_csm_agent()",
         "复制 create_ops_agent() → 改 schema / policy / router / executors\n"
         "audit_logger / session_store / mcp_client 等都继续用 Kernel factory"),
    ]
    y = 1.6
    for num, title, file, code in steps:
        # num
        c = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(y + 0.1), Inches(0.55), Inches(0.55)
        )
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
        c.line.fill.background()
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        set_font(r, size=16, bold=True, color=WHITE)

        # title + file
        add_textbox(s, Inches(1.3), Inches(y), Inches(11), Inches(0.35),
                    title, size=14, bold=True, color=NAVY)
        add_textbox(s, Inches(1.3), Inches(y + 0.32), Inches(11), Inches(0.3),
                    file, size=11, color=ACCENT, font=MONO_FONT)

        # code-ish body
        add_rect(s, Inches(1.3), Inches(y + 0.65), Inches(11.5), Inches(0.35 * max(3, code.count("\n") + 1)),
                 fill=CODE_BG)
        nlines = code.split("\n")
        for i, line in enumerate(nlines):
            add_textbox(s, Inches(1.45), Inches(y + 0.7 + i * 0.25),
                        Inches(11.3), Inches(0.25),
                        line, size=10, color=CODE_TEXT, font=MONO_FONT)
        y += 1.05

    add_footer(s)


# ====================== Part 7  Evolution ======================


def slide_29_degradation(prs):
    s = blank_slide(prs)
    add_header(s, "三级降级路径（§10）",
               "每一级都有明确触发条件和用户感知",
               page_num=29, chapter="PART 7  演进方向与测试")

    levels = [
        ("L1", "Executor 级降级",
         LIGHT_BLUE, DEEP_BLUE,
         "触发：单个 executor 抛异常",
         "处理：PlanStepStatus=FAILED → fail-fast → 返回错误说明",
         "用户看到：\"步骤 X 执行失败：xxx\"；其他步骤不继续"),
        ("L2", "Planner 级降级",
         RGBColor(0xFF, 0xF2, 0xE6), ACCENT,
         "触发：max_iterations 耗尽 / Planner 生成空 Plan",
         "处理：fallback_plan 兜底 → 单 knowledge step",
         "用户看到：AI 进入「普通问答」模式回答"),
        ("L3", "Kernel 级降级",
         RGBColor(0xFF, 0xE8, 0xE0), RED,
         "触发：Receipt 失败 / 记忆 Backend 故障 / 整个 graph 崩",
         "处理：chat() 外层 try/except → ChatResponse 不崩",
         "用户看到：\"系统暂时不可用\"；Audit 仍然落盘错误条目"),
    ]
    y = 1.6
    for code, title, bg, fg, trigger, handle, ux in levels:
        add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.6), fill=bg)
        # 左侧大 tag
        tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.7), Inches(y + 0.15),
                                  Inches(1.5), Inches(1.3))
        tag.fill.solid(); tag.fill.fore_color.rgb = fg
        tag.line.fill.background()
        tf = tag.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = code
        set_font(r, size=32, bold=True, color=WHITE, font=ENGLISH_FONT)

        add_textbox(s, Inches(2.4), Inches(y + 0.12), Inches(10), Inches(0.4),
                    title, size=17, bold=True, color=fg)
        add_multiline(s, Inches(2.4), Inches(y + 0.55), Inches(10), Inches(1.05),
                      [trigger, handle, ux], size=12, color=DARK_GRAY,
                      line_spacing=1.25)
        y += 1.75


def slide_30_testing(prs):
    s = blank_slide(prs)
    add_header(s, "测试金字塔与目前覆盖情况",
               "L0 单元 / L1 契约 / L2 E2E —— 共 84 个自动化测试用例",
               page_num=30, chapter="PART 7  演进方向与测试")

    # Pyramid
    # L2 (top, narrow)
    l2 = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                            Inches(4.5), Inches(1.55),
                            Inches(4.3), Inches(1.3))
    l2.fill.solid(); l2.fill.fore_color.rgb = PURPLE
    l2.line.fill.background()
    add_textbox(s, Inches(4.5), Inches(2.2), Inches(4.3), Inches(0.5),
                "L2 · E2E\n27 个", size=13, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # L1 middle
    l1 = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID,
                            Inches(3.5), Inches(2.85),
                            Inches(6.3), Inches(1.2))
    l1.fill.solid(); l1.fill.fore_color.rgb = DEEP_BLUE
    l1.line.fill.background()
    add_textbox(s, Inches(3.5), Inches(3.1), Inches(6.3), Inches(0.8),
                "L1 · Kernel 契约\n4 个 · 保证 Kernel 不与 Ops 耦合",
                size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # L0 bottom
    l0 = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID,
                            Inches(2.0), Inches(4.05),
                            Inches(9.3), Inches(1.6))
    l0.fill.solid(); l0.fill.fore_color.rgb = GREEN
    l0.line.fill.background()
    add_textbox(s, Inches(2.0), Inches(4.35), Inches(9.3), Inches(1.0),
                "L0 · 单元测试\n53 个 · Planner / Registry / Topology / Memory / Patterns ...",
                size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Right panel: details
    add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.15),
             fill=LIGHT_GRAY)
    add_textbox(s, Inches(0.6), Inches(6.08), Inches(12), Inches(0.35),
                "E2E 测试矩阵（27 个用例分 6 组）",
                size=13, bold=True, color=NAVY)
    groups = [
        ("A", "Happy Path", 5, GREEN),
        ("B", "Kernel 不变量", 8, RED),
        ("C", "插件点", 7, PURPLE),
        ("D", "Vertical 隔离", 2, DEEP_BLUE),
        ("E", "降级路径", 2, ACCENT),
        ("F", "反模式回归", 3, MEDIUM_GRAY),
    ]
    x = 0.6
    for code, name, cnt, color in groups:
        add_rect_text(s, Inches(x), Inches(6.5), Inches(2.0), Inches(0.5),
                      f"{code}  {name}  {cnt}", fill=color, text_color=WHITE,
                      size=11, bold=True)
        x += 2.05


def slide_31_supervisor(prs):
    s = blank_slide(prs)
    add_header(s, "未来演进：Supervisor 多 Agent 协同",
               "跨域问题拆给多个 Vertical，最后汇总 —— §7 演进方向",
               page_num=31, chapter="PART 7  演进方向与测试")

    # Query
    add_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.55),
             fill=NAVY)
    add_textbox(s, Inches(0.6), Inches(1.55), Inches(12.1), Inches(0.4),
                "用户：\"Q3 订单为什么下滑？\"",
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Supervisor
    add_rect_text(s, Inches(4.5), Inches(2.2), Inches(4.3), Inches(0.8),
                  "Supervisor\n（MetaPlanner + AgentProxyExecutor）",
                  fill=ACCENT, text_color=WHITE, size=13, bold=True)
    # arrows down
    for x in [2.1, 5.0, 7.9, 10.8]:
        add_arrow(s, 6.65, 3.05, x + 1.0, 3.6, color=MEDIUM_GRAY, width=1)

    # 4 子 Agent
    sub = [
        ("DataAgent", "查销售曲线\nSQL + BI", GREEN),
        ("CsmAgent", "投诉分类\n退款 / 延迟占比", DEEP_BLUE),
        ("JARVIS", "线上异常排查\n影响转化的故障", ACCENT),
        ("DocAgent", "生成一页\n摘要报告", PURPLE),
    ]
    x = 1.1
    for name, body, color in sub:
        add_rect_text(s, Inches(x), Inches(3.6), Inches(2.8), Inches(1.3),
                      f"{name}\n\n{body}", fill=color, text_color=WHITE,
                      size=12, bold=True)
        # arrow down
        add_arrow(s, x + 1.4, 4.95, x + 1.4, 5.45, color=MEDIUM_GRAY, width=1)
        x += 2.9

    # Aggregator
    add_rect_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.7),
                  "Supervisor 聚合：联合审批 / 跨 Agent Audit / 统一 ChatResponse",
                  fill=NAVY, text_color=WHITE, size=14, bold=True)

    # key mechanism
    add_rect(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.75),
             fill=LIGHT_BLUE)
    add_textbox(s, Inches(0.7), Inches(6.48), Inches(12), Inches(0.3),
                "🔑 关键机制：PlanStep.execution_target = \"agent:data\" / \"agent:csm\" / ...",
                size=12, bold=True, color=NAVY, font=MONO_FONT)
    add_textbox(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.3),
                "已在 PlanStep Schema 预留字段 — 不需要重构 Kernel 就能落地 Supervisor",
                size=11, color=DARK_GRAY)


def slide_32_roadmap(prs):
    s = blank_slide(prs)
    add_header(s, "总结 · 路线图 · Q&A",
               "", page_num=32, chapter="PART 7  演进方向与测试")

    # 左：总结
    add_rect(s, Inches(0.5), Inches(1.5), Inches(6.15), Inches(5.3),
             fill=LIGHT_BLUE, line=DEEP_BLUE)
    add_textbox(s, Inches(0.7), Inches(1.6), Inches(5.8), Inches(0.4),
                "✅ 本次讲解覆盖", size=16, bold=True, color=DEEP_BLUE)
    summary = [
        "• Kernel 11 个组件 + 5 条不变量",
        "• 10 个插件点 + 新 Vertical 五步清单",
        "• JARVIS 4 个 Executor + 中文拆分",
        "• 审批凭据 / 多假设诊断 详细流程",
        "• 三级降级 + 84 个自动化测试覆盖",
        "• Supervisor 演进方向已预留字段",
    ]
    y = 2.15
    for line in summary:
        add_textbox(s, Inches(0.9), Inches(y), Inches(5.5), Inches(0.35),
                    line, size=13, color=DARK_GRAY)
        y += 0.4
    add_rect(s, Inches(0.5), Inches(5.3), Inches(6.15), Inches(1.45),
             fill=NAVY)
    add_textbox(s, Inches(0.7), Inches(5.4), Inches(5.8), Inches(0.4),
                "核心价值", size=14, bold=True, color=ACCENT)
    add_textbox(s, Inches(0.7), Inches(5.8), Inches(5.8), Inches(0.9),
                "JARVIS 只是起点 ——\n我们在造一个可以孵化任意\n业务 Agent 的基建平台",
                size=13, bold=True, color=WHITE)

    # 右：路线图
    add_rect(s, Inches(6.85), Inches(1.5), Inches(6.0), Inches(5.3),
             fill=RGBColor(0xFF, 0xF2, 0xE6), line=ACCENT)
    add_textbox(s, Inches(7.05), Inches(1.6), Inches(5.7), Inches(0.4),
                "🛣️  路线图", size=16, bold=True, color=ACCENT)

    stages = [
        ("近期", "Q2",
         ["DiagnosisExecutor 接入 MultiHypothesisExecutor 基类",
          "接入真实 MCP 服务器（k8s-mcp / jenkins-mcp）",
          "补齐 RedisSessionStore / MemoryBackend"]),
        ("中期", "Q3-Q4",
         ["落地第二个 Vertical（客服 / 数据 任选）",
          "Kernel 跨域通用性验证",
          "Agent 的灰度 / 回滚 / 版本管理"]),
        ("远期", "1 年+",
         ["Supervisor 多 Agent 协同落地",
          "跨域请求自动拆解",
          "统一人机交互与观测平台"]),
    ]
    y = 2.1
    for stage, period, items in stages:
        add_rect_text(s, Inches(7.05), Inches(y), Inches(1.0), Inches(0.5),
                      stage, fill=ACCENT, text_color=WHITE, size=12, bold=True)
        add_textbox(s, Inches(8.1), Inches(y + 0.05), Inches(4.6), Inches(0.4),
                    period, size=11, color=MEDIUM_GRAY, font=ENGLISH_FONT)
        for i, item in enumerate(items):
            add_textbox(s, Inches(7.25), Inches(y + 0.55 + i * 0.3),
                        Inches(5.5), Inches(0.3),
                        f"▸ {item}", size=11, color=DARK_GRAY)
        y += 1.55

    add_footer(s)


# ---------- main ----------


SLIDES = [
    slide_01_cover,
    slide_02_toc,
    slide_03_problem,
    slide_04_split_insight,
    slide_05_principles,
    slide_06_layers_overview,
    slide_07_boundary_rules,
    slide_08_directory,
    slide_09_kernel_components,
    slide_10_base_agent_graph,
    slide_11_planner,
    slide_12_router,
    slide_13_executor,
    slide_14_tools,
    slide_15_memory,
    slide_16_approval,
    slide_17_audit,
    slide_18_invariants,
    slide_19_ops_composition,
    slide_20_ops_executors,
    slide_21_ops_planner,
    slide_22_ops_approval_memory,
    slide_23_chat_flow,
    slide_24_approval_gate,
    slide_25_diagnosis,
    slide_26_compound_example,
    slide_27_plugin_points,
    slide_28_new_vertical,
    slide_29_degradation,
    slide_30_testing,
    slide_31_supervisor,
    slide_32_roadmap,
]


def main():
    prs = new_presentation()
    for func in SLIDES:
        func(prs)
    out_dir = Path("docs")
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "ops_agent_architecture_briefing.pptx"
    prs.save(out_path)
    print(f"✅ Generated: {out_path.resolve()}")
    print(f"   Slides:    {len(prs.slides)}")


if __name__ == "__main__":
    main()
